#!/usr/bin/env python3
"""
Skill Builder Web App
Run: python3 app.py
Open: http://localhost:5173
"""

import io
import json
import os
import queue
import re
import shutil
import sys
import tempfile
import threading
import uuid
import zipfile
from pathlib import Path

from flask import Flask, Response, jsonify, render_template, request, send_file, stream_with_context

# Use bundled scripts/ dir (works locally and on Railway)
SKILL_SCRIPTS = Path(__file__).parent / "scripts"
sys.path.insert(0, str(SKILL_SCRIPTS))

app = Flask(__name__)
jobs = {}  # job_id -> {status, queue, result, zip_path, skill_dir, skill_name, error}

# Drafts directory — persists within a Railway session (survives browser refresh/disconnect)
# Uses Railway volume if mounted, otherwise /tmp
DRAFTS_DIR = Path(os.getenv("RAILWAY_VOLUME_MOUNT_PATH", "/tmp")) / "skillbuilder_drafts"
DRAFTS_DIR.mkdir(parents=True, exist_ok=True)


# ──────────────────────────────────────────────
# Background build worker
# ──────────────────────────────────────────────

BATCH_SIZE = 35_000          # chars per Opus batch
CONTENT_BUDGET = 60_000      # max chars of summarized content sent to Opus
SUMMARY_MAX_INPUT = 15_000   # max chars of raw content fed to Haiku per source
MAX_PRE_SUMMARIZE = 300      # hard cap on sources sent to Haiku (ranked by title relevance first)

# Haiku pricing: $0.80/M input, $4/M output (approx)
# Opus pricing:  $15/M input, $75/M output (approx)
HAIKU_INPUT_COST_PER_CHAR  = 0.80  / 1_000_000
HAIKU_OUTPUT_COST_PER_CHAR = 4.00  / 1_000_000
OPUS_INPUT_COST_PER_CHAR   = 15.0  / 1_000_000
OPUS_OUTPUT_COST_PER_CHAR  = 75.0  / 1_000_000
SUMMARY_OUTPUT_CHARS = 1_500   # expected output chars per Haiku summary

SUMMARIZE_PROMPT = """\
You are extracting the most valuable knowledge from source material for a skill-building system.

Extract ONLY the actionable, specific, reusable knowledge — tactics, frameworks, formulas,
step-by-step processes, specific numbers, rules of thumb, and decision criteria.

Strip out: stories, filler, repetition, motivational content, and anything generic.

Format as tight bullet points grouped by topic. Target 250-350 words.
Do NOT add commentary or preamble — just the bullets.
"""

def _parse_skill_json(raw):
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        m = re.search(r"\{[\s\S]+\}", raw)
        if m:
            try:
                return json.loads(m.group(0))
            except Exception:
                pass
    return None


def _make_batches(all_parts):
    """Split content parts into batches of ~BATCH_SIZE chars each."""
    batches = []
    current_batch = []
    current_size = 0
    for part in all_parts:
        content = part["content"]
        # If a single part exceeds batch size, split it
        while len(content) > BATCH_SIZE:
            chunk = content[:BATCH_SIZE]
            current_batch.append({**part, "content": chunk})
            batches.append(current_batch)
            current_batch = []
            current_size = 0
            content = content[BATCH_SIZE:]
        if current_size + len(content) > BATCH_SIZE and current_batch:
            batches.append(current_batch)
            current_batch = []
            current_size = 0
        current_batch.append({**part, "content": content})
        current_size += len(content)
    if current_batch:
        batches.append(current_batch)
    return batches


def _summarize_content(client, content, title, source_intention, global_intention):
    """Summarize a single piece of content using Haiku. Returns summarized text."""
    # Per-source intention takes priority; fall back to global
    effective_intention = source_intention or global_intention
    intent_line = f"\nFocus on content relevant to: {effective_intention}\n" if effective_intention else ""
    try:
        msg = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=600,
            system=SUMMARIZE_PROMPT,
            messages=[{"role": "user", "content":
                f"Source: {title}{intent_line}\n\n{content[:SUMMARY_MAX_INPUT]}"
            }],
        )
        return msg.content[0].text.strip()
    except Exception:
        return content[:SUMMARY_OUTPUT_CHARS]


def _score_summary(summary, intention):
    """Score a summary by keyword density against the intention."""
    if not intention:
        return 1
    import re as _re
    stop = {"a","an","the","and","or","but","in","on","at","to","for","of","with",
            "by","from","as","is","are","was","were","be","been","have","has","had",
            "do","does","did","will","would","could","should","i","we","you","they","it"}
    keywords = [w for w in _re.findall(r'\b[a-z]{3,}\b', intention.lower()) if w not in stop]
    if not keywords:
        return 1
    text = summary.lower()
    return sum(text.count(kw) for kw in keywords)


def _title_score(title, intention):
    """Quick keyword relevance score on title alone (no Haiku needed)."""
    if not intention:
        return 1
    import re as _re
    stop = {"a","an","the","and","or","but","in","on","at","to","for","of","with",
            "by","from","as","is","are","was","were","be","been","have","has","had",
            "do","does","did","will","would","could","should","i","we","you","they","it"}
    keywords = [w for w in _re.findall(r'\b[a-z]{3,}\b', intention.lower()) if w not in stop]
    if not keywords:
        return 1
    text = title.lower()
    return sum(text.count(kw) for kw in keywords)


def run_build(job_id, sources, max_videos, raw_text, intention=""):
    """
    Pipeline:
    1. Fetch all content (free)
    2. Pre-rank by title relevance, cap at MAX_PRE_SUMMARIZE
    3. Cost estimate gate (before Haiku spend) → user confirms or cancels
    4. Haiku summarizes the capped set
    5. Re-rank summaries, fill CONTENT_BUDGET
    6. Opus builds skill from budget content
    """
    q = jobs[job_id]["queue"]

    def log(msg, kind="info"):
        q.put({"type": kind, "message": msg})

    try:
        from fetch_content import fetch
        from generate_skill import SYSTEM_PROMPT, build_user_message, read_knowledge_base, save_skill
        import anthropic

        # source_intentions: {url: intention_string}
        source_intentions = jobs[job_id].get("source_intentions", {})

        all_parts = []

        # ── Phase 1: Fetch all content ────────────────────────────────────────
        for url in sources:
            if jobs[job_id].get("cancelled"):
                return
            log(f"Fetching: {url}")
            src_intention = source_intentions.get(url, "")
            effective_intention = src_intention or intention
            try:
                def _fwd(msg):
                    if jobs[job_id].get("cancelled"):
                        return
                    q.put(msg)
                result = fetch(url, max_videos=max_videos, verbose=False,
                               intention=effective_intention, log_fn=_fwd)

                # For collections, split into per-video parts for granular ranking
                if result.get("videos"):
                    for v in result["videos"]:
                        all_parts.append({
                            "title": v["title"],
                            "content": v["transcript"],
                            "url": v["url"],
                            "source_url": url,
                            "source_intention": src_intention,
                        })
                else:
                    all_parts.append({
                        "title": result["title"],
                        "content": result["content"],
                        "url": url,
                        "source_url": url,
                        "source_intention": src_intention,
                    })

                vids = result.get("video_count", 1)
                chars = len(result["content"])
                label = f"{vids} videos, {chars:,} chars" if vids > 1 else f"{chars:,} chars"
                log(f"Got: {result['title']} ({label})", "success")
            except Exception as e:
                log(f"Failed: {url} — {e}", "error")

        if raw_text:
            all_parts.append({
                "title": "Uploaded Content",
                "content": raw_text,
                "url": "",
                "source_url": "",
                "source_intention": "",
            })
            log(f"Got: Uploaded content ({len(raw_text):,} chars)", "success")

        if not all_parts:
            raise ValueError("No content was successfully fetched. Check URLs and try again.")

        if jobs[job_id].get("cancelled"):
            return

        # ── Phase 2: Pre-rank by title relevance, cap at MAX_PRE_SUMMARIZE ───
        total_fetched = len(all_parts)
        if len(all_parts) > MAX_PRE_SUMMARIZE:
            # Score by title keywords first (free) so we keep the most relevant
            for part in all_parts:
                part["title_score"] = _title_score(part["title"], intention)
            all_parts.sort(key=lambda p: -p["title_score"])
            all_parts = all_parts[:MAX_PRE_SUMMARIZE]
            log(f"Large source set: kept top {MAX_PRE_SUMMARIZE} of {total_fetched} sources by title relevance", "info")

        # ── Phase 3: Cost estimate gate (BEFORE spending on Haiku) ───────────
        total_raw = sum(min(len(p["content"]), SUMMARY_MAX_INPUT) for p in all_parts)
        haiku_input_chars  = total_raw
        haiku_output_chars = len(all_parts) * SUMMARY_OUTPUT_CHARS
        haiku_cost = (haiku_input_chars * HAIKU_INPUT_COST_PER_CHAR +
                      haiku_output_chars * HAIKU_OUTPUT_COST_PER_CHAR)

        # Estimate Opus cost from budget (worst case: full CONTENT_BUDGET)
        est_budget = min(len(all_parts) * SUMMARY_OUTPUT_CHARS, CONTENT_BUDGET)
        est_batches = max(1, est_budget // BATCH_SIZE + 1)
        opus_input_chars  = est_budget + 10_000
        opus_output_chars = est_batches * 6_000
        opus_cost = (opus_input_chars * OPUS_INPUT_COST_PER_CHAR +
                     opus_output_chars * OPUS_OUTPUT_COST_PER_CHAR)

        total_cost = haiku_cost + opus_cost
        est_secs   = len(all_parts) * 3 + est_batches * 45 + 30

        q.put({
            "type": "cost_estimate",
            "sources": len(all_parts),
            "total_sources": total_fetched,
            "batches": est_batches,
            "haiku_cost": round(haiku_cost, 4),
            "opus_cost": round(opus_cost, 4),
            "total_cost": round(total_cost, 4),
            "estimated_secs": est_secs,
        })

        if total_cost > 1.00:
            jobs[job_id]["status"] = "awaiting_confirmation"
            jobs[job_id]["pending_parts"] = all_parts
            jobs[job_id]["pending_source_url"] = sources[0] if sources else ""
            return  # Frontend must POST /confirm/<job_id> to continue

        # Under $1 — proceed automatically
        _run_summarize_and_opus(job_id, all_parts, sources, intention, q, log)

    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        log(f"Error: {e}", "error")
        q.put({"type": "complete", "skill_name": None})


def _run_summarize_and_opus(job_id, all_parts, sources, intention, q, log):
    """Phases 4-6: Haiku summarize → rank → fill budget → Opus build."""
    import anthropic
    client = anthropic.Anthropic()

    try:
        # ── Phase 4: Haiku summarizes the capped set ──────────────────────────
        log(f"Summarizing {len(all_parts)} sources with Haiku...")
        q.put({"type": "summarize_total", "count": len(all_parts)})

        for i, part in enumerate(all_parts, 1):
            if jobs[job_id].get("cancelled"):
                return
            summary = _summarize_content(
                client, part["content"], part["title"],
                part.get("source_intention", ""), intention
            )
            part["summary"] = summary
            q.put({"type": "summarize_progress", "current": i, "total": len(all_parts)})

        log(f"Summarization complete — {len(all_parts)} sources distilled", "success")

        # ── Phase 5: Re-rank summaries, fill content budget ───────────────────
        for part in all_parts:
            part["score"] = _score_summary(part["summary"], intention)

        ranked = sorted(all_parts, key=lambda p: -p["score"])

        selected = []
        budget_used = 0
        for part in ranked:
            s = part["summary"]
            if budget_used + len(s) <= CONTENT_BUDGET:
                selected.append(part)
                budget_used += len(s)
            if budget_used >= CONTENT_BUDGET:
                break

        skipped = len(all_parts) - len(selected)
        log(f"Selected top {len(selected)} sources ({budget_used:,} chars) — {skipped} lower-relevance excluded", "info")

        batches = _make_batches([{**p, "content": p["summary"]} for p in selected])

        # ── Phase 6: Opus ─────────────────────────────────────────────────────
        _run_opus_phase(job_id, selected, batches, sources, intention, q, log, client)

    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        log(f"Error: {e}", "error")
        q.put({"type": "complete", "skill_name": None})


def _save_draft(job_id, partial_skills, intention, source_title, n_batches, batch_num):
    """Write current partial skill state to disk so it survives browser disconnects."""
    try:
        draft_dir = DRAFTS_DIR / job_id
        draft_dir.mkdir(parents=True, exist_ok=True)
        # Merge what we have so far for preview
        if partial_skills:
            combined_md = "\n\n---\n\n".join(p.get("skill_md", "") for p in partial_skills)
            skill_name = partial_skills[-1].get("skill_name", "Draft Skill")
            draft = {
                "job_id": job_id,
                "skill_name": skill_name,
                "skill_md": combined_md,
                "summary": partial_skills[-1].get("summary", ""),
                "intention": intention,
                "source_title": source_title,
                "batches_done": batch_num,
                "batches_total": n_batches,
                "is_partial": batch_num < n_batches,
            }
            (draft_dir / "partial.json").write_text(json.dumps(draft, ensure_ascii=False))
    except Exception:
        pass  # Draft saving is best-effort — never crash the build


def _run_opus_phase(job_id, selected, batches, sources, intention, q, log, client):
    """Phase 6: Opus builds skill from ranked summaries."""
    from generate_skill import SYSTEM_PROMPT, build_user_message, read_knowledge_base

    source_title = "Combined: " + ", ".join(p["title"][:20] for p in selected[:3])
    source_url   = sources[0] if sources else ""
    n_batches    = len(batches)

    jobs[job_id]["status"] = "running"

    log("Loading knowledge base...")
    best_practices, lessons, sources_catalog = read_knowledge_base()

    if n_batches > 1:
        est_secs = n_batches * 45 + 30
        q.put({"type": "batch_total", "count": n_batches, "estimated_secs": est_secs})

    partial_skills = []

    for i, batch in enumerate(batches, 1):
        if jobs[job_id].get("cancelled"):
            return

        batch_text = "\n\n---\n\n".join(
            f"# {p['title']}\n\n{p['summary']}" for p in batch
        )
        log(f"Generating batch {i}/{n_batches} with Claude Opus 4.6...")

        try:
            msg = client.messages.create(
                model="claude-opus-4-6",
                max_tokens=6000,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": build_user_message(
                    batch_text, source_title, source_url,
                    best_practices, lessons, sources_catalog,
                    intention=intention,
                )}],
            )
            raw = msg.content[0].text.strip()
            result = _parse_skill_json(raw)

            if not result:
                log(f"Batch {i} parse failed — retrying...", "info")
                msg2 = client.messages.create(
                    model="claude-opus-4-6",
                    max_tokens=6000,
                    system=SYSTEM_PROMPT,
                    messages=[{"role": "user", "content": build_user_message(
                        batch_text, source_title, source_url,
                        best_practices, lessons, sources_catalog,
                        intention=intention,
                    )}],
                )
                result = _parse_skill_json(msg2.content[0].text.strip())

            if result:
                partial_skills.append(result)
                # Save draft to disk after every completed batch
                _save_draft(job_id, partial_skills, intention, source_title, n_batches, i)
                # Also expose partial result in-memory for polling
                jobs[job_id]["partial_skill"] = {
                    "skill_name": result.get("skill_name", "Draft"),
                    "skill_md": "\n\n---\n\n".join(p.get("skill_md","") for p in partial_skills),
                    "batches_done": i,
                    "batches_total": n_batches,
                }
                q.put({"type": "batch_done", "batch": i, "total": n_batches,
                       "skill_name": result.get("skill_name", "Draft")})
                log(f"Batch {i}/{n_batches} complete ✓", "success")
            else:
                log(f"Batch {i}/{n_batches} failed after retry — skipping", "error")

        except Exception as e:
            log(f"Batch {i} error: {e}", "error")

        if n_batches > 1:
            q.put({"type": "batch_progress", "current": i, "total": n_batches})

    if not partial_skills:
        jobs[job_id]["status"] = "error"
        q.put({"type": "complete", "skill_name": None})
        return

    # Merge if needed
    if len(partial_skills) == 1:
        skill_result = partial_skills[0]
    else:
        log(f"Merging {len(partial_skills)} batches into final skill...", "info")
        MERGE_PROMPT = """\
Merge these partial Claude Code skill files into one final coherent skill.
1. Synthesize all knowledge — keep the best version of each point
2. Remove redundancy
3. Keep under 500 lines
4. Return JSON with keys: skill_name, skill_md, sources_md, notes_md, summary, questions
"""
        partials_text = "\n\n---PARTIAL---\n\n".join(p.get("skill_md","") for p in partial_skills)
        merge_msg = client.messages.create(
            model="claude-opus-4-6",
            max_tokens=8000,
            system=MERGE_PROMPT,
            messages=[{"role": "user", "content":
                f"Merge {len(partial_skills)} partials. Intention: {intention}\n\n{partials_text}"
            }],
        )
        skill_result = _parse_skill_json(merge_msg.content[0].text.strip())
        if not skill_result:
            log("Merge parse failed — using best single batch", "error")
            skill_result = partial_skills[0]
        else:
            log("Merge complete ✓", "success")

    questions = skill_result.get("questions", [])
    if questions:
        log(f"Found {len(questions)} conflict(s) that need your input", "warn")
        jobs[job_id].update({
            "status": "needs_input",
            "result": skill_result,
            "questions": questions,
            "combined_content": "\n\n".join(p.get("skill_md","") for p in partial_skills),
            "source_title": source_title,
            "source_url": source_url,
        })
        q.put({"type": "needs_input", "questions": questions})
        return

    _finalize_skill(job_id, skill_result, q)


def _write_watchlist(target_dir: Path, sources: list, intention: str, target_type: str, name: str):
    """Write watchlist.json and initial CHANGELOG.md if any sources have a schedule."""
    from datetime import date as _date
    recurring = [s for s in sources if s.get("schedule") and s["schedule"] != "never"]
    watchlist = {
        "name": name,
        "type": target_type,
        "intention": intention,
        "sources": [
            {
                "url": s["url"],
                "schedule": s.get("schedule", "never"),
                "auto_accept": False,
                "last_fetched": _date.today().isoformat(),
            }
            for s in sources
        ],
        "created": _date.today().isoformat(),
        "last_updated": _date.today().isoformat(),
    }
    (target_dir / "watchlist.json").write_text(
        json.dumps(watchlist, indent=2), encoding="utf-8"
    )
    # Initial changelog
    cl = target_dir / "CHANGELOG.md"
    source_lines = "\n".join(
        f"- {s['url']}" + (f" ({s['schedule']})" if s.get("schedule") and s["schedule"] != "never" else "")
        for s in sources
    )
    cl.write_text(
        f"# {name} — Changelog\n\n"
        f"## {_date.today().isoformat()} (created)\n"
        f"Initial build from:\n{source_lines}\n",
        encoding="utf-8",
    )


def _finalize_skill(job_id, skill_result, q):
    """Save skill to disk, build zip, mark done."""
    from generate_skill import save_skill
    tmpdir = Path(tempfile.mkdtemp())
    skill_dir = save_skill(skill_result, tmpdir)

    # Write watchlist.json + CHANGELOG.md if recurring sources exist
    job = jobs[job_id]
    watch_sources = job.get("watch_sources", [])
    intention = job.get("intention", "")
    skill_name = skill_result["skill_name"]
    if watch_sources:
        _write_watchlist(skill_dir, watch_sources, intention, "skill", skill_name)

    zip_path = tmpdir / f"{skill_name}.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for file in skill_dir.rglob("*"):
            if file.is_file():
                zf.write(file, file.relative_to(tmpdir))

    jobs[job_id].update({
        "status": "done",
        "result": skill_result,
        "zip_path": str(zip_path),
        "skill_dir": str(skill_dir),
        "skill_name": skill_name,
    })

    # Remove draft file — build is complete
    try:
        draft_dir = DRAFTS_DIR / job_id
        if draft_dir.exists():
            shutil.rmtree(draft_dir)
    except Exception:
        pass

    q = jobs[job_id]["queue"]
    q.put({"type": "complete", "skill_name": skill_name})


# ──────────────────────────────────────────────
# Persona background worker
# ──────────────────────────────────────────────

def run_build_persona(job_id, sources, max_videos, raw_text, intention=""):
    q = jobs[job_id]["queue"]

    def log(msg, kind="info"):
        q.put({"type": kind, "message": msg})

    try:
        from fetch_content import fetch
        from generate_persona import generate_persona

        all_parts = []

        for url in sources:
            if jobs[job_id].get("cancelled"):
                return
            log(f"Fetching: {url}")
            try:
                def _fwd_p(msg):
                    if jobs[job_id].get("cancelled"):
                        return
                    q.put(msg)
                result = fetch(url, max_videos=max_videos, verbose=False,
                               intention=intention, log_fn=_fwd_p)
                all_parts.append({
                    "title": result["title"],
                    "content": result["content"],
                    "url": url,
                    "video_count": result.get("video_count", 1),
                })
                vids = result.get("video_count", 1)
                chars = len(result["content"])
                label = f"{vids} videos, {chars:,} chars" if vids > 1 else f"{chars:,} chars"
                log(f"Got: {result['title']} ({label})", "success")
            except Exception as e:
                log(f"Failed: {url} — {e}", "error")

        if raw_text:
            all_parts.append({
                "title": "Uploaded Content",
                "content": raw_text,
                "url": "",
                "video_count": 1,
            })
            log(f"Got: Uploaded content ({len(raw_text):,} chars)", "success")

        if not all_parts:
            raise ValueError("No content was successfully fetched. Check URLs and try again.")

        if jobs[job_id].get("cancelled"):
            return

        if len(all_parts) == 1:
            combined = all_parts[0]["content"]
            source_title = all_parts[0]["title"]
            source_url = all_parts[0]["url"]
        else:
            parts_text = []
            for p in all_parts:
                header = f"# {p['title']}"
                if p["url"]:
                    header += f"\nSource: {p['url']}"
                parts_text.append(f"{header}\n\n{p['content']}")
            combined = "\n\n---\n\n".join(parts_text)
            source_title = "Combined: " + ", ".join(p["title"][:25] for p in all_parts[:3])
            source_url = sources[0] if sources else ""

        MAX_P = 40_000
        if len(combined) > MAX_P:
            log(f"Content truncated to {MAX_P:,} chars for processing ({len(combined):,} total)")
            combined = combined[:MAX_P] + "\n\n[...truncated...]"

        log("Generating persona with Claude Opus 4.6...")
        result = generate_persona(combined, source_title, source_url, intention=intention)

        persona_name = result["persona_name"]

        # Save persona to ~/.claude/personas/<slug>/
        import re as _re
        slug = _re.sub(r"[^a-z0-9]+", "-", persona_name.lower()).strip("-")
        persona_root = Path.home() / ".claude" / "personas" / slug
        persona_root.mkdir(parents=True, exist_ok=True)
        persona_file = persona_root / f"{slug}.md"
        persona_file.write_text(result["persona_md"], encoding="utf-8")

        # Write watchlist + changelog if recurring sources
        watch_sources = jobs[job_id].get("watch_sources", [])
        intention_val = jobs[job_id].get("intention", "")
        if watch_sources:
            _write_watchlist(persona_root, watch_sources, intention_val, "persona", slug)
        else:
            # Always write a basic changelog
            from datetime import date as _date
            cl = persona_root / "CHANGELOG.md"
            cl.write_text(
                f"# {persona_name} — Changelog\n\n"
                f"## {_date.today().isoformat()} (created)\n"
                f"Initial build from: {source_url or source_title}\n",
                encoding="utf-8",
            )

        jobs[job_id].update({
            "status": "done",
            "result": result,
            "skill_name": persona_name,
            "persona_dir": str(persona_root),
        })
        q.put({"type": "complete", "skill_name": persona_name})

    except Exception as e:
        jobs[job_id]["status"] = "error"
        jobs[job_id]["error"] = str(e)
        log(f"Error: {e}", "error")
        q.put({"type": "complete", "skill_name": None})


# ──────────────────────────────────────────────
# Routes
# ──────────────────────────────────────────────

@app.route("/")
def landing():
    return render_template("landing.html")


@app.route("/app")
def index():
    return render_template("index.html")


@app.route("/build", methods=["POST"])
def build():
    if not os.getenv("ANTHROPIC_API_KEY"):
        return jsonify({"error": "ANTHROPIC_API_KEY is not set. Export it in your terminal before starting this app."}), 400

    data = request.json
    sources = [s.strip() for s in data.get("sources", []) if s.strip()]
    raw_text = data.get("raw_text", "").strip()
    intention = data.get("intention", "").strip()
    max_videos = int(data.get("max_videos", 50))
    files = data.get("files", [])  # [{name, content}]
    # watch_sources: [{url, schedule}] — sources with recurring schedules
    watch_sources = data.get("watch_sources", [])
    # source_intentions: {url: intention_string} — per-source focus
    source_intentions = data.get("source_intentions", {})

    # Merge uploaded file content into raw_text
    if files:
        file_blocks = [f"# {f['name']}\n\n{f['content']}" for f in files]
        if raw_text:
            file_blocks.append(raw_text)
        raw_text = "\n\n---\n\n".join(file_blocks)

    if not sources and not raw_text:
        return jsonify({"error": "Add at least one URL, file, or paste some text."}), 400

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "status": "running",
        "queue": queue.Queue(),
        "log_buffer": [],
        "result": None, "zip_path": None,
        "skill_dir": None, "skill_name": None, "error": None,
        "cancelled": False,
        "watch_sources": watch_sources,
        "intention": intention,
        "sources": sources,
        "source_intentions": source_intentions,
    }

    t = threading.Thread(target=run_build, args=(job_id, sources, max_videos, raw_text, intention))
    t.daemon = True
    t.start()

    return jsonify({"job_id": job_id})


@app.route("/build-persona", methods=["POST"])
def build_persona_route():
    if not os.getenv("ANTHROPIC_API_KEY"):
        return jsonify({"error": "ANTHROPIC_API_KEY is not set. Export it in your terminal before starting this app."}), 400

    data = request.json
    sources = [s.strip() for s in data.get("sources", []) if s.strip()]
    raw_text = data.get("raw_text", "").strip()
    intention = data.get("intention", "").strip()
    max_videos = int(data.get("max_videos", 50))
    files = data.get("files", [])
    watch_sources = data.get("watch_sources", [])

    if files:
        file_blocks = [f"# {f['name']}\n\n{f['content']}" for f in files]
        if raw_text:
            file_blocks.append(raw_text)
        raw_text = "\n\n---\n\n".join(file_blocks)

    if not sources and not raw_text:
        return jsonify({"error": "Add at least one URL, file, or paste some text."}), 400

    job_id = str(uuid.uuid4())[:8]
    jobs[job_id] = {
        "status": "running",
        "queue": queue.Queue(),
        "log_buffer": [],
        "result": None, "zip_path": None,
        "skill_dir": None, "skill_name": None, "error": None,
        "job_type": "persona",
        "cancelled": False,
        "watch_sources": watch_sources,
        "intention": intention,
    }

    t = threading.Thread(target=run_build_persona, args=(job_id, sources, max_videos, raw_text, intention))
    t.daemon = True
    t.start()

    return jsonify({"job_id": job_id})


@app.route("/stream/<job_id>")
def stream(job_id):
    if job_id not in jobs:
        return "Not found", 404

    def generate():
        q = jobs[job_id]["queue"]
        buf = jobs[job_id]["log_buffer"]
        while True:
            try:
                msg = q.get(timeout=15)
                buf.append(msg)
                yield f"data: {json.dumps(msg)}\n\n"
                if msg.get("type") == "complete":
                    break
            except queue.Empty:
                yield "data: {\"type\":\"ping\"}\n\n"

    return Response(
        stream_with_context(generate()),
        mimetype="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@app.route("/poll/<job_id>")
def poll(job_id):
    """Polling fallback: returns buffered log messages since index N."""
    job = jobs.get(job_id)
    if not job:
        # Check disk drafts — job may have been from a previous session
        draft_file = DRAFTS_DIR / job_id / "partial.json"
        if draft_file.exists():
            return jsonify({"error": "Job not in memory", "has_draft": True, "job_id": job_id}), 404
        return jsonify({"error": "Job not found"}), 404
    since = int(request.args.get("since", 0))
    buf = job.get("log_buffer", [])
    return jsonify({
        "messages": buf[since:],
        "total": len(buf),
        "status": job.get("status"),
        "partial_skill": job.get("partial_skill"),  # included when Opus batches complete
    })


@app.route("/drafts")
def list_drafts():
    """List all saved draft builds."""
    drafts = []
    for draft_dir in sorted(DRAFTS_DIR.iterdir(), key=lambda p: p.stat().st_mtime, reverse=True):
        partial = draft_dir / "partial.json"
        if partial.exists():
            try:
                data = json.loads(partial.read_text())
                drafts.append({
                    "job_id": data.get("job_id", draft_dir.name),
                    "skill_name": data.get("skill_name", "Untitled"),
                    "intention": data.get("intention", ""),
                    "batches_done": data.get("batches_done", 0),
                    "batches_total": data.get("batches_total", 1),
                    "is_partial": data.get("is_partial", False),
                    "saved_at": int(partial.stat().st_mtime * 1000),
                })
            except Exception:
                pass
    return jsonify({"drafts": drafts})


@app.route("/draft/<job_id>")
def get_draft(job_id):
    """Return the saved draft skill for a job."""
    draft_file = DRAFTS_DIR / job_id / "partial.json"
    if not draft_file.exists():
        return jsonify({"error": "No draft found"}), 404
    try:
        data = json.loads(draft_file.read_text())
        return jsonify(data)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/draft/<job_id>", methods=["DELETE"])
def delete_draft(job_id):
    """Delete a saved draft."""
    draft_dir = DRAFTS_DIR / job_id
    if draft_dir.exists():
        shutil.rmtree(draft_dir)
    return jsonify({"ok": True})


@app.route("/result/<job_id>")
def result(job_id):
    job = jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    r = job.get("result") or {}
    is_persona = job.get("job_type") == "persona"
    return jsonify({
        "status": job["status"],
        "job_type": job.get("job_type", "skill"),
        "skill_name": job.get("skill_name"),
        "skill_md": r.get("skill_md") if not is_persona else None,
        "sources_md": r.get("sources_md") if not is_persona else None,
        "persona_name": r.get("persona_name") if is_persona else None,
        "persona_md": r.get("persona_md") if is_persona else None,
        "topics_covered": r.get("topics_covered") if is_persona else None,
        "summary": r.get("summary"),
        "questions": job.get("questions", []),
        "error": job.get("error"),
    })


@app.route("/answer/<job_id>", methods=["POST"])
def answer(job_id):
    """Receive user answers to contradiction questions, regenerate skill with context."""
    job = jobs.get(job_id)
    if not job or job["status"] != "needs_input":
        return jsonify({"error": "Job not in needs_input state"}), 400

    answers = request.json.get("answers", [])  # [{question, answer}]

    # Build answers context string
    answers_text = "\n".join(
        f"Q: {a['question']}\nA: {a['answer']}" for a in answers
    )

    # Re-run generation with answers injected
    job["status"] = "running"
    job["questions"] = []

    def regenerate():
        q = job["queue"]
        try:
            from generate_skill import SYSTEM_PROMPT, build_user_message, read_knowledge_base
            import anthropic

            q.put({"type": "info", "message": "Regenerating skill with your answers..."})
            best_practices, lessons, sources_catalog = read_knowledge_base()

            # Inject answers into the user message
            base_msg = build_user_message(
                job["combined_content"], job["source_title"], job["source_url"],
                best_practices, lessons, sources_catalog,
            )
            msg_with_answers = base_msg + f"\n\n## User Answers to Contradictions\n{answers_text}\n\nApply these answers when resolving the conflicts you identified."

            client = anthropic.Anthropic()
            message = client.messages.create(
                model="claude-opus-4-6",
                max_tokens=6000,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": msg_with_answers}],
            )
            raw = message.content[0].text.strip()
            try:
                skill_result = json.loads(raw)
            except json.JSONDecodeError:
                m = re.search(r"\{[\s\S]+\}", raw)
                skill_result = json.loads(m.group(0)) if m else None
            if not skill_result:
                raise ValueError("Could not parse regenerated skill")

            job["result"] = skill_result
            q.put({"type": "success", "message": f"Skill updated with your answers"})
            _finalize_skill(job_id, skill_result, q)

        except Exception as e:
            job["status"] = "error"
            job["error"] = str(e)
            q.put({"type": "error", "message": str(e)})
            q.put({"type": "complete", "skill_name": None})

    t = threading.Thread(target=regenerate)
    t.daemon = True
    t.start()

    return jsonify({"ok": True})


@app.route("/append/<job_id>", methods=["POST"])
def append_source(job_id):
    """Fetch a new URL and merge it into an existing skill."""
    job = jobs.get(job_id)
    if not job or job["status"] != "done" or job.get("job_type") == "persona":
        return jsonify({"error": "No completed skill found for this session."}), 400

    data = request.json
    url = data.get("url", "").strip()
    max_videos = int(data.get("max_videos", 50))
    if not url:
        return jsonify({"error": "No URL provided."}), 400

    try:
        from fetch_content import fetch

        result = fetch(url, max_videos=max_videos, verbose=False)
        new_content = result["content"]
        source_title = result["title"]

        MAX = 40_000
        if len(new_content) > MAX:
            new_content = new_content[:MAX] + "\n\n[...truncated...]"

        existing_skill_md = job["result"]["skill_md"]

        import anthropic, re as _re
        APPEND_PROMPT = """\
You are updating an existing Claude Code skill with new content.

Your job:
1. Identify genuinely new information, techniques, or updated recommendations in the new content
2. Update the existing skill to incorporate it — revise sections, add bullets, update examples
3. Do NOT bloat the skill — remove outdated information if the new content supersedes it
4. Keep the skill under 400 lines
5. Return the complete updated SKILL.md text only — no JSON, no preamble

If the new content adds nothing meaningful, return the existing skill unchanged.
"""
        client = anthropic.Anthropic()
        msg = client.messages.create(
            model="claude-opus-4-6",
            max_tokens=6000,
            system=APPEND_PROMPT,
            messages=[{"role": "user", "content":
                f"## Existing SKILL.md\n\n{existing_skill_md}\n\n"
                f"---\n\n## New Content from: {source_title}\nURL: {url}\n\n{new_content}"
            }],
        )
        updated_skill_md = msg.content[0].text.strip()

        # Update in-memory result
        job["result"]["skill_md"] = updated_skill_md

        # Rebuild zip with updated skill
        from generate_skill import save_skill
        tmpdir = Path(tempfile.mkdtemp())
        skill_result = dict(job["result"])
        skill_dir = save_skill(skill_result, tmpdir)
        zip_path = tmpdir / f"{job['skill_name']}.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for file in skill_dir.rglob("*"):
                if file.is_file():
                    zf.write(file, file.relative_to(tmpdir))
        job["zip_path"] = str(zip_path)
        job["skill_dir"] = str(skill_dir)

        return jsonify({"skill_md": updated_skill_md})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/cancel/<job_id>", methods=["POST"])
def cancel(job_id):
    job = jobs.get(job_id)
    if job:
        job["cancelled"] = True
        job["status"] = "cancelled"
    return jsonify({"ok": True})


@app.route("/confirm/<job_id>", methods=["POST"])
def confirm_build(job_id):
    """Resume an awaiting_confirmation job after user approves the cost estimate."""
    job = jobs.get(job_id)
    if not job:
        return jsonify({"error": "Job not found"}), 404
    if job.get("status") != "awaiting_confirmation":
        return jsonify({"error": "Job is not awaiting confirmation"}), 400

    all_parts = job.pop("pending_parts", [])
    sources   = job.get("sources", [])
    intention = job.get("intention", "")
    q         = job["queue"]

    job["status"] = "running"

    def _log(m, k="info"):
        q.put({"type": k, "message": m})

    def _resume():
        _run_summarize_and_opus(job_id, all_parts, sources, intention, q, _log)

    import threading
    t = threading.Thread(target=_resume, daemon=True)
    t.start()
    return jsonify({"ok": True})


@app.route("/download/<job_id>/<file_type>")
def download(job_id, file_type):
    job = jobs.get(job_id)
    if not job or job["status"] != "done":
        return "Not ready", 404

    if file_type == "zip":
        return send_file(
            job["zip_path"], as_attachment=True,
            download_name=f"{job['skill_name']}.zip",
        )
    elif file_type == "skill_md":
        content = job["result"]["skill_md"].encode("utf-8")
        return send_file(
            io.BytesIO(content), as_attachment=True,
            download_name="SKILL.md", mimetype="text/markdown",
        )
    elif file_type == "sources_md":
        content = job["result"]["sources_md"].encode("utf-8")
        return send_file(
            io.BytesIO(content), as_attachment=True,
            download_name="sources.md", mimetype="text/markdown",
        )
    elif file_type == "persona_md":
        content = job["result"]["persona_md"].encode("utf-8")
        name = (job["result"].get("persona_name") or "persona").replace(" ", "-") + ".md"
        return send_file(
            io.BytesIO(content), as_attachment=True,
            download_name=name, mimetype="text/markdown",
        )
    return "Unknown type", 404


@app.route("/parse-file", methods=["POST"])
def parse_file():
    """Accept a binary file upload and return extracted plain text."""
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    filename = f.filename or ""
    ext = Path(filename).suffix.lower()

    try:
        if ext in (".txt", ".md"):
            text = f.read().decode("utf-8", errors="ignore")

        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(f)
            text = "\n".join(page.extract_text() or "" for page in reader.pages)

        elif ext in (".docx",):
            from docx import Document
            doc = Document(f)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext in (".pptx",):
            from pptx import Presentation
            prs = Presentation(f)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
            text = "\n\n".join(parts)

        elif ext in (".xlsx", ".xls"):
            import openpyxl, anthropic as _anthropic
            wb = openpyxl.load_workbook(f, data_only=False)
            sections = []
            all_formulas = []

            for sheet in wb.worksheets:
                rows_text = []
                sheet_formulas = []

                for row in sheet.iter_rows():
                    row_vals = []
                    for cell in row:
                        val = cell.value
                        if val is None:
                            row_vals.append("")
                        elif isinstance(val, str) and val.startswith("="):
                            placeholder = f"[FORMULA:{cell.coordinate}]"
                            row_vals.append(placeholder)
                            sheet_formulas.append({
                                "cell": cell.coordinate,
                                "formula": val,
                                "sheet": sheet.title,
                            })
                        else:
                            row_vals.append(str(val))
                    # Skip fully empty rows
                    if any(v for v in row_vals):
                        rows_text.append(" | ".join(row_vals))

                section = f"## Sheet: {sheet.title}\n" + "\n".join(rows_text)
                sections.append(section)
                all_formulas.extend(sheet_formulas)

            raw_text = "\n\n".join(sections)

            # Ask Claude to explain all formulas in one call
            if all_formulas:
                formula_list = "\n".join(
                    f"- {f['sheet']}!{f['cell']}: {f['formula']}" for f in all_formulas
                )
                client = _anthropic.Anthropic()
                msg = client.messages.create(
                    model="claude-opus-4-6",
                    max_tokens=2000,
                    messages=[{"role": "user", "content":
                        f"Explain each of these spreadsheet formulas in plain English. "
                        f"Be specific about what each one calculates and why it's useful. "
                        f"Format as a list with the cell reference, then the explanation.\n\n{formula_list}"
                    }],
                )
                formula_explanations = msg.content[0].text.strip()
                text = raw_text + f"\n\n## Formula Explanations\n\n{formula_explanations}"
            else:
                text = raw_text

        else:
            return jsonify({"error": f"Unsupported file type: {ext}"}), 400

        return jsonify({"name": filename, "content": text, "chars": len(text)})

    except Exception as e:
        return jsonify({"error": f"Could not parse {filename}: {e}"}), 500


@app.route("/install/<job_id>", methods=["POST"])
def install(job_id):
    job = jobs.get(job_id)
    if not job or job["status"] != "done":
        return jsonify({"error": "Job not ready"}), 400

    src = Path(job["skill_dir"])
    dest = Path.home() / ".claude" / "skills" / job["skill_name"]
    if dest.exists():
        shutil.rmtree(dest)
    shutil.copytree(src, dest)

    return jsonify({
        "success": True,
        "installed_to": str(dest),
        "message": f"Installed to {dest}. Restart Claude Code to activate.",
    })


if __name__ == "__main__":
    port = int(os.getenv("PORT", 5173))
    print(f"\n  Skill Builder → http://localhost:{port}\n")
    if not os.getenv("ANTHROPIC_API_KEY"):
        print("  WARNING: ANTHROPIC_API_KEY not set\n")
    app.run(host="0.0.0.0", debug=False, port=port, threaded=True)
